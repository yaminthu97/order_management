from pptx import Presentation
import sys;
import os
import json

def replace_text_in_presentation(pptx_path, json_load):
    # PowerPointファイルを読み込む
    prs = Presentation(pptx_path)

    # 各スライドをループ
    for slide in prs.slides:
        # 各シェイプをループ
        for shape in slide.shapes:
            # シェイプがテキストフレームを持っている場合
            if shape.has_text_frame:
                # 各段落をループ
                for paragraph in shape.text_frame.paragraphs:
                    # 各ランをループ
                    for run in paragraph.runs:
                        if(run.text in json_load.keys()):
                        # 検索文字列を置き換える
                            if(json_load[run.text]):
                                run.text = run.text.replace(run.text, json_load[run.text])
                            else:
                                run.text = run.text.replace(run.text, "")

    # 変更を保存する
    prs.save(sys.argv[3])


if(len(sys.argv) != 4):
	print("usage: python "+sys.argv[0] + " templateFile jsonFile outputFile")
	sys.exit(101)

#fファイル有無(テンプレートファイル)
if not os.path.isfile(sys.argv[1]):
	print("テンプレートファイルがありません")
	sys.exit(102)

#fファイル有無(jsonファイル)
if not os.path.isfile(sys.argv[2]):
	print("JSONファイルがありません")
	sys.exit(103)

try:
	json_open = open(sys.argv[2], 'r')
	json_load = json.load(json_open)
except Exception as e:
	print("JSONファイルの読み込みに失敗しました")
	sys.exit(104)

#pptxチェック
try:
    # PowerPointファイルを読み込む
    Presentation(sys.argv[1])
except Exception as e:
	print("pptxファイルの読み込みに失敗しました")
	sys.exit(105)

replace_text_in_presentation(sys.argv[1], json_load)
